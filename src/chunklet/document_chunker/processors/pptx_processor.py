from collections.abc import Generator
from typing import Any

# pptx is lazily imported
from chunklet.document_chunker.md_table import build_md_table
from chunklet.document_chunker.processors.base_processor import BaseProcessor


class PPTXProcessor(BaseProcessor):
    """
    Processor class for extracting text, tables, charts, notes, and metadata from PPTX files.

    Text content is extracted sequentially slide-by-slide. Structural elements
    like slide titles are converted to Markdown headers, bullet points maintain
    indentation, presentation tables are structured into valid Markdown tables,
    visual charts are transformed into text grids, and presenter notes are appended
    at the bottom of each slide block.

    This processor focuses on extracting core **metadata** following the OpenXML
    Document CoreProperties format, matching common practice in office document types.

    For more details on PPTX layout elements, refer to the `python-pptx` documentation:
    https://python-pptx.readthedocs.io/
    """

    # PowerPoint built-in structural placeholder IDs to filter out:
    # 14 = HEADER, 15 = FOOTER, 16 = SLIDE_NUMBER
    TEMPLATE_PLACEHOLDERS = {14, 15, 16}

    METADATA_FIELDS = [
        "title",
        "author",
        "last_modified_by",
        "created",
        "modified",
        "subject",
        "keywords",
        "category",
        "content_status",
        "version",
    ]

    def __init__(self, file_path: str):
        """
        Initializes the PPTXProcessor with a path to the PPTX file
        and reads the presentation into memory.

        Args:
            file_path: Path to the PPTX file.
        """
        super().__init__(file_path)
        try:
            from pptx import Presentation
        except ImportError as e:  # pragma: no cover
            raise ImportError(
                "The 'python-pptx' library is not installed. "
                "Please install it with 'pip install python-pptx>=1.0.0' or install the document processing extras "
                "with 'pip install 'chunklet-py[structured-document]''"
            ) from e

        self.prs = Presentation(file_path)

    def _extract_slide_title(self, slide: Any) -> str | None:
        """
        Safely isolates and extracts the authentic presentation slide title shape.

        Args:
            slide: A python-pptx slide object.

        Returns:
            The string text of the title formatted as a Markdown H1 header, or None.
        """
        title_shape = getattr(slide.shapes, "title", None)
        if (
            title_shape
            and hasattr(title_shape, "has_text_frame")
            and title_shape.has_text_frame
        ):
            title_text = title_shape.text.strip()
            if title_text and not title_text.isdigit():
                return f"# {title_text}\n"
        return None

    def _extract_text(self, shape: Any) -> str | None:
        """
        Processes standard shapes, placeholders, and content text frames.

        Args:
            shape: A python-pptx shape object.

        Returns:
            Extracted lines of text formatted into Markdown paragraphs or bullets.
        """
        if not shape.has_text_frame:
            return None

        if self._is_template_placeholder(shape):
            return None

        lines = []
        for p in shape.text_frame.paragraphs:
            text = p.text.strip()
            if not text or text.isdigit():
                continue
            prefix = "" if p.level == 0 else "  " * (p.level - 1) + "- "
            lines.append(f"{prefix}{text}")

        return "\n".join(lines) if lines else None

    def _is_template_placeholder(self, shape: Any) -> bool:
        """Check if a shape is a background template element to skip."""
        if not shape.is_placeholder:
            return False
        pf = getattr(shape, "placeholder_format", None)
        return getattr(pf, "type", None) in self.TEMPLATE_PLACEHOLDERS

    def _extract_table(self, shape: Any) -> str | None:
        """
        Extracts complex layout grid shapes and transforms them into Markdown tables.

        Args:
            shape: A python-pptx shape object holding table attributes.

        Returns:
            A clean structured Markdown grid table text block.
        """
        if not hasattr(shape, "has_table") or not shape.has_table:
            return None

        table = shape.table
        rows = [
            [cell.text.strip().replace("\n", " ") for cell in row.cells]
            for row in table.rows
        ]
        if not rows:
            return None

        return build_md_table(rows)

    def _extract_chart(self, shape: Any) -> str | None:
        """
        Extracts chart properties, categories, and plot series values into Markdown.

        Args:
            shape: A python-pptx graphic frame shape object containing a chart.

        Returns:
            A structured Markdown representation of the chart data, or None.
        """
        if not hasattr(shape, "has_chart") or not shape.has_chart:
            return None

        chart = shape.chart
        parts = ["\n### [Chart]"]

        title = self._safe_chart_title(chart)
        if title:
            parts.append(f"**Title**: {title}")

        for plot in getattr(chart, "plots", []):
            table = self._plot_to_table(plot)
            if table:
                parts.append("")
                parts.append(build_md_table(table))

        parts.append("")
        return "\n".join(parts)

    @staticmethod
    def _safe_chart_title(chart: Any) -> str | None:
        """Extract chart title text safely, avoiding missing-attribute errors."""
        title = getattr(chart, "chart_title", None)
        if not title or not hasattr(title, "has_text_frame"):
            return None
        if not title.has_text_frame:
            return None
        text = title.text_frame.text.strip()
        return text or None

    @staticmethod
    def _plot_to_table(plot: Any) -> list[list[str]] | None:
        """Build a list-of-lists table (header + rows) from a chart plot."""
        if not (categories := getattr(plot, "categories", None)):
            return None
        categories = [str(c) for c in categories]

        series_list = getattr(plot, "series", [])
        if not series_list:
            return None

        headers = ["Categories"] + [
            getattr(s, "name", f"Series {i + 1}") or f"Series {i + 1}"
            for i, s in enumerate(series_list)
        ]
        rows = [headers]
        for cat_idx, cat in enumerate(categories):
            row = [cat]
            for s in series_list:
                values = getattr(s, "values", [])
                try:
                    val = values[cat_idx] if cat_idx < len(values) else ""
                    row.append(str(val) if val is not None else "")
                except (IndexError, AttributeError, TypeError):
                    row.append("")
            rows.append(row)
        return rows

    def _extract_notes(self, slide: Any) -> str | None:
        """
        Extracts presenter/speaker notes text associated with the slide.

        Args:
            slide: A python-pptx slide object.

        Returns:
            Extracted text blocks formatted as a blockquote block, or None.
        """
        notes_slide = getattr(slide, "notes_slide", None)
        if (
            notes_slide
            and hasattr(notes_slide, "notes_text_frame")
            and notes_slide.notes_text_frame
        ):
            notes_lines = []
            for paragraph in notes_slide.notes_text_frame.paragraphs:
                text = paragraph.text.strip()
                if not text:
                    continue

                level = paragraph.level
                if level == 0:
                    notes_lines.append(f"> {text}\n>")
                else:
                    indent = "  " * (level - 1)
                    notes_lines.append(f"> {indent}- {text}")

            if notes_lines:
                return "\n***\n**Notes:**\n" + "\n".join(notes_lines).rstrip(">")
        return None

    def extract_metadata(self) -> dict[str, Any]:
        """
        Extracts OpenXML Document CoreProperties from the PPTX file
        based on the defined METADATA_FIELDS class schema.

        Returns:
            A dictionary containing metadata fields.
        """
        meta = self.prs.core_properties
        metadata = {"source": str(self.file_path)}

        for field in self.METADATA_FIELDS:
            # Handle cases where property might not exist in an older python-pptx build
            if not hasattr(meta, field):
                continue

            val = getattr(meta, field, None)
            if val:
                metadata[field] = str(val)

        return metadata

    def extract_text(self) -> Generator[str, None, None]:
        """
        Yields fully converted Markdown content slide-by-slide from the PPTX archive.

        Yields:
            Markdown-formatted rendering of each slide.
        """
        for slide_idx, slide in enumerate(self.prs.slides, start=1):
            slide_content = [f"\n<!-- Slide {slide_idx} -->\n"]

            # 1. Title Processing
            title = self._extract_slide_title(slide)
            if title:
                slide_content.append(title)

            # 2. Iterate remaining layout blocks
            for shape in slide.shapes:
                if shape == slide.shapes.title:
                    continue

                # Handle Table blocks
                if hasattr(shape, "has_table") and shape.has_table:
                    table_md = self._extract_table(shape)
                    if table_md:
                        slide_content.append(table_md)

                # Handle Chart blocks
                elif hasattr(shape, "has_chart") and shape.has_chart:
                    chart_md = self._extract_chart(shape)
                    if chart_md:
                        slide_content.append(chart_md)

                # Handle text boxes & paragraphs
                elif hasattr(shape, "has_text_frame") and shape.has_text_frame:
                    text_md = self._extract_text(shape)
                    if text_md:
                        slide_content.append(text_md)

            # 3. Presenter Notes Append
            notes_md = self._extract_notes(slide)
            if notes_md:
                slide_content.append(notes_md)

            # Package slide buffer string cleanly
            yield "\n".join(slide_content).strip()


# --- Example usage ---
if __name__ == "__main__":  # pragma: no cover
    file_path = "samples/presentation.pptx"
    processor = PPTXProcessor(file_path)

    metadata = processor.extract_metadata()
    print("Metadata:")
    for k, v in metadata.items():
        print(f"{k}: {v}")

    print("\nSlide Content Stream:\n")
    for _, slide_markdown in enumerate(processor.extract_text(), start=1):
        print(slide_markdown)
